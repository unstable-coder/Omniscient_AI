from __future__ import annotations
from typing import Iterable
from docx import Document
from pptx import Presentation
from app.parsers.base_parser import BaseParser
from app.models.schemas import ContentUnit

class DocumentParser(BaseParser):
    supported_extensions = ("docx", "pptx")
    supported_mime_types = (
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    def parse(self, path: str, metadata: dict[str, object]) -> Iterable[ContentUnit]:
        suffix = path.rsplit('.', 1)[-1].lower()
        if suffix == "docx":
            yield from self._parse_docx(path, metadata)
        elif suffix == "pptx":
            yield from self._parse_pptx(path, metadata)

    def _parse_docx(self, path: str, metadata: dict[str, object]) -> Iterable[ContentUnit]:
        doc = Document(path)
        for paragraph_index, paragraph in enumerate(doc.paragraphs, start=1):
            text = paragraph.text.strip()
            if text:
                unit_metadata = {**metadata, "paragraph_index": paragraph_index}
                yield ContentUnit(text=text, metadata=unit_metadata)
        for table_index, table in enumerate(doc.tables, start=1):
            rows = []
            for row in table.rows:
                rows.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
            if rows:
                unit_metadata = {**metadata, "table_index": table_index}
                yield ContentUnit(text="\n".join(rows), metadata=unit_metadata)

    def _parse_pptx(self, path: str, metadata: dict[str, object]) -> Iterable[ContentUnit]:
        presentation = Presentation(path)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text.strip())
                elif shape.has_table:
                    for row in shape.table.rows:
                        slide_text.append(" | ".join(cell.text.strip() for cell in row.cells if cell.text.strip()))
            text = "\n".join([line for line in slide_text if line])
            if text:
                unit_metadata = {**metadata, "slide_number": slide_index}
                yield ContentUnit(text=text, metadata=unit_metadata)
